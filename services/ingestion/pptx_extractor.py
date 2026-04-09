"""PPTX ingestion service for CloudInsight."""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from services.ingestion.common import (
    IngestionValidationError,
    build_error_result,
    build_success_result,
    clean_text,
    count_words,
    join_text_fragments,
    validate_file,
)


logger = logging.getLogger(__name__)


class PPTXExtractor:
    """Extract slide-wise text from PowerPoint presentations."""

    def extract(self, file_path: str | Path) -> dict[str, Any]:
        """Extract presentation text grouped by slide."""
        try:
            path = validate_file(file_path, {".pptx"})
        except IngestionValidationError as exc:
            return build_error_result(file_path, "pptx", errors=[str(exc)])

        try:
            from pptx import Presentation
        except ImportError:
            return build_error_result(
                path,
                "pptx",
                errors=["Missing dependency 'python-pptx'. Install it before ingesting PPTX files."],
            )

        try:
            presentation = Presentation(str(path))
            sections: list[dict[str, Any]] = []
            empty_slide_count = 0

            for slide_number, slide in enumerate(presentation.slides, start=1):
                fragments = self._extract_slide_fragments(slide)
                slide_text = join_text_fragments(fragments)
                title = clean_text(slide.shapes.title.text) if getattr(slide.shapes, "title", None) else ""

                if not slide_text:
                    empty_slide_count += 1
                    continue

                sections.append(
                    {
                        "id": f"slide_{slide_number}",
                        "type": "slide",
                        "slide_number": slide_number,
                        "title": title,
                        "text": slide_text,
                        "char_count": len(slide_text),
                        "word_count": count_words(slide_text),
                    }
                )

            warnings = []
            if empty_slide_count:
                warnings.append(f"{empty_slide_count} slide(s) contained no extractable text.")

            metadata = {
                "slide_count": len(presentation.slides),
                "extracted_slide_count": len(sections),
                "empty_slide_count": empty_slide_count,
            }
            full_text = join_text_fragments(section["text"] for section in sections)

            logger.info("Extracted PPTX content from %s", path)
            return build_success_result(
                path,
                "pptx",
                metadata=metadata,
                text=full_text,
                sections=sections,
                warnings=warnings,
            )
        except Exception as exc:
            logger.exception("PPTX extraction failed for %s", path)
            return build_error_result(path, "pptx", errors=[f"PPTX extraction failed: {exc}"])

    def _extract_slide_fragments(self, slide: Any) -> list[str]:
        """Collect text fragments from all supported shapes within a slide."""
        fragments: list[str] = []
        for shape in slide.shapes:
            fragments.extend(self._extract_shape_fragments(shape))
        return fragments

    def _extract_shape_fragments(self, shape: Any) -> list[str]:
        """Recursively collect text from text, table, and grouped shapes."""
        fragments: list[str] = []

        if hasattr(shape, "shapes"):
            for child in shape.shapes:
                fragments.extend(self._extract_shape_fragments(child))

        if getattr(shape, "has_text_frame", False):
            text = clean_text(shape.text or "")
            if text:
                fragments.append(text)

        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                row_text = " | ".join(clean_text(cell.text) for cell in row.cells if clean_text(cell.text))
                if row_text:
                    fragments.append(row_text)

        return fragments


def extract_pptx(file_path: str | Path, **_: Any) -> dict[str, Any]:
    """Convenience wrapper for PPTX extraction."""
    return PPTXExtractor().extract(file_path)
